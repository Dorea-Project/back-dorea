"""La validation du livrable — **l'ordre, qui est toute la protection**.

Le cœur pur (`test_livrable.py`) sait juger deux chaînes. Ce fichier garde ce qui l'entoure, et
qui décide vraiment :

1. **Rien n'est produit sans une division du plan.** C'est la règle centrale, et elle est
   arithmétique : le document met en page ce que le pasteur a écrit.
2. **Le jugement précède le fichier.** Une diapositive altérée rend `rejete` — et il n'y a rien
   à supprimer, puisque rien n'a été produit. Un contrôle d'après coup protège la base de
   données, pas l'assemblée.
3. **`conforme` est signé.** Par celui qui valide, pas par l'auteur de la préparation : c'est
   lui qui répond de ce qui sortira.
4. **Toutes les versions détenues sont consultées** (Q9), et la version reconnue est enregistrée
   — sur Romains 8:1 elle change la doctrine du verset projeté.
5. **Une référence qui n'existe pas ne bloque que sa diapositive**, avec le motif du corpus : on
   rend le dossier entier, pas la première faute.
"""

from __future__ import annotations

from datetime import UTC, datetime
from uuid import uuid4

import pytest

from app.contexts.urim.application.ports import ElementRecord, PreparationRecord
from app.contexts.urim.deliverable.application.ports import (
    DiapositiveSoumise,
    TexteServi,
)
from app.contexts.urim.deliverable.application.service import UrimDeliverableService
from app.contexts.urim.deliverable.domain.citation import ALTERE, EXACT
from app.contexts.urim.deliverable.domain.documents import POINT_CENTRAL
from app.contexts.urim.domain.errors import (
    LivrableNonValideError,
    LivrableSansPlanError,
)

from .test_study_service import AUTEUR, EGLISE, MAINTENANT, UNITE, _Acces, _index

pytestmark = pytest.mark.asyncio

#: Hébreux 13:1 tel que l'index de test le sert.
HB_13_1 = "Persévérez dans l'amour fraternel."
#: La même phrase avec un mot changé — l'altération que le contrôle existe pour voir.
HB_13_1_ALTERE = "Persévérez dans l'amour du prochain."
#: Une seconde version détenue, qui porte une formulation différente et légitime.
HB_13_1_AUTRE = "Que l'amour fraternel demeure."

LSG, AUTRE_VERSION = uuid4(), uuid4()


class _Studies:
    def __init__(self, record: PreparationRecord | None, elements=()) -> None:
        self._record = record
        self._elements = list(elements)

    async def get(self, study_id):
        if self._record is not None and study_id == self._record.id:
            return self._record
        return None

    async def list_elements(self, study_id):
        return self._elements


class _Livrables:
    def __init__(self) -> None:
        self.ecrits: list[tuple] = []

    async def add(self, record, controles) -> None:
        self.ecrits.append((record, controles))

    async def get(self, deliverable_id):
        return next((r for r, _ in self.ecrits if r.id == deliverable_id), None)

    async def controles(self, deliverable_id):
        return next((c for r, c in self.ecrits if r.id == deliverable_id), [])


class _Versets:
    """Deux versions détenues — c'est le minimum pour que Q9 veuille dire quelque chose."""

    def __init__(self, *textes: TexteServi) -> None:
        self._textes = list(textes) or [
            TexteServi(LSG, "LSG", HB_13_1),
            TexteServi(AUTRE_VERSION, "Ostervald", HB_13_1_AUTRE),
        ]
        self.appels: list[tuple] = []

    async def textes(
        self, *, book_id, chapter, verse_start, verse_end, prefer_version_id=None
    ):
        self.appels.append((book_id, chapter, verse_start, verse_end))
        return list(self._textes)


def _preparation(**kw) -> PreparationRecord:
    return PreparationRecord(
        id=kw.pop("id", uuid4()),
        church_id=kw.pop("church_id", EGLISE),
        author_id=kw.pop("author_id", AUTEUR),
        raw_input="Hébreux 13:1-2",
        pericope_id=UNITE,
        version_id=LSG,
        **kw,
    )


class _Etude:
    """Le dossier rejoué. Il **note qu'on l'a lu** — le rendu de la note en dépend."""

    def __init__(self, dto=None) -> None:
        self.dto = dto
        self.lectures: list = []

    async def get(self, *, actor_account_id, study_id):
        self.lectures.append(study_id)
        return self.dto


def _dossier_d_etude(prep: PreparationRecord):
    """Un `StudyDTO` minimal — **avec ce que la note doit imprimer et le deck refuser**."""
    from types import SimpleNamespace

    return SimpleNamespace(
        record=prep,
        resolved_label="Hébreux 13:1-2",
        pericope_label="Exhortations",
        pericope_reviewed_by="ia-mistral",
        trace=[("bound_pericope", "L'unité tient du v. 1 au v. 2.")],
        elements=[ElementRecord(POINT_CENTRAL, 1, "1- L'amour fraternel demeure")],
        verses=[SimpleNamespace(reference="Hébreux 13:1", text=HB_13_1)],
        bearings=[
            SimpleNamespace(
                axis_code="ecclesiologie", label="ecclésiologie",
                strength="dominant", rationale="le texte exhorte l'assemblée",
            )
        ],
        caveats=["Le texte ne dit pas quand l'Esprit est donné."],
        couples=[
            SimpleNamespace(
                plan_source="thematique", subject_matter="doctrine",
                feasible=False, refusal_reason="aucun personnage",
                proof_text_risk="eleve",
            )
        ],
        resisting_elsewhere=[
            SimpleNamespace(label="2 Co 12:7-10", rationale="Dieu dit non trois fois")
        ],
        supports=[("Hb 13v1", "Hébreux 13:1", HB_13_1, "exact")],
        options=[("opt-2", "La péricope entière", "écartée le 12/08", "curation", True)],
    )


def _service(record, *, elements=(), livrables=None, versets=None, etude=None):
    return UrimDeliverableService(
        studies=_Studies(record, elements),
        etude=etude or _Etude(),
        livrables=livrables or _Livrables(),
        versets=versets or _Versets(),
        access=_Acces(),
        index=_index(),
        clock=lambda: MAINTENANT,
    )


_PLAN = (ElementRecord(POINT_CENTRAL, 1, "1- L'amour fraternel demeure"),)


# ============================================================ 1. rien sans son plan


async def test_sans_division_du_plan_aucun_livrable_n_est_produit():
    """**La règle centrale**, et elle est arithmétique : sans plan, il n'y a rien à imprimer.

    Le motif oriente — un refus qui n'oriente pas est une porte fermée."""
    prep = _preparation()
    depot = _Livrables()
    service = _service(prep, elements=(), livrables=depot)

    with pytest.raises(LivrableSansPlanError) as refus:
        await service.soumettre(actor_account_id=AUTEUR, study_id=prep.id)

    assert "ne l'écrit pas à votre place" in str(refus.value)
    assert depot.ecrits == []  # rien n'a été écrit, pas même un rejet


async def test_un_titre_seul_ne_suffit_pas():
    """Le couple : une étiquette n'est pas un plan."""
    prep = _preparation()
    service = _service(prep, elements=(ElementRecord("titre", 1, "L'amour fraternel"),))
    with pytest.raises(LivrableSansPlanError):
        await service.soumettre(actor_account_id=AUTEUR, study_id=prep.id)


async def test_une_division_suffit_a_ouvrir_le_livrable():
    prep = _preparation()
    service = _service(prep, elements=_PLAN)
    dto = await service.soumettre(actor_account_id=AUTEUR, study_id=prep.id, kind="note")
    assert dto.conforme
    assert dto.record.format == "docx"  # la note a son format natif


# ============================================================ 2. le jugement précède le fichier


async def test_une_diapositive_alteree_rend_rejete_et_aucun_fichier_n_existe():
    """**L'ordre est toute la protection.** Le service ne produit aucun octet : il valide.

    Un contrôle d'après coup protégerait la base de données, pas l'assemblée — un fichier
    produit est un fichier qui circule."""
    prep = _preparation()
    service = _service(prep, elements=_PLAN)

    dto = await service.soumettre(
        actor_account_id=AUTEUR,
        study_id=prep.id,
        diapositives=[DiapositiveSoumise("Le texte", "Hébreux 13:1", HB_13_1_ALTERE)],
    )

    assert dto.record.validation == "rejete"
    assert dto.controles[0].verdict == ALTERE
    # …et personne n'a signé un rejet.
    assert dto.record.validated_by is None and dto.record.validated_at is None


async def test_un_texte_fidele_rend_conforme_et_signe():
    """`conforme` est **signé**, et par celui qui valide : c'est lui qui monte en chaire."""
    prep = _preparation(author_id=uuid4())  # la préparation est d'un collègue
    service = _service(prep, elements=_PLAN)

    dto = await service.soumettre(
        actor_account_id=AUTEUR,
        study_id=prep.id,
        diapositives=[DiapositiveSoumise("Le texte", "Hébreux 13:1", HB_13_1)],
    )

    assert dto.record.validation == "conforme"
    assert dto.record.validated_by == AUTEUR
    assert dto.record.validated_at == MAINTENANT


# ============================================================ 3. toutes les versions (Q9)


async def test_le_texte_d_une_autre_version_detenue_est_reconnu_et_enregistre():
    """Q9 de bout en bout : un pasteur cite la Bible qu'il a.

    La version reconnue descend dans `citation_check.version_id` — la colonne existait et rien
    ne la remplissait."""
    prep = _preparation()
    service = _service(prep, elements=_PLAN)

    dto = await service.soumettre(
        actor_account_id=AUTEUR,
        study_id=prep.id,
        diapositives=[DiapositiveSoumise("Le texte", "Hébreux 13:1", HB_13_1_AUTRE)],
    )

    assert dto.conforme
    assert dto.controles[0].verdict == EXACT
    assert dto.controles[0].version_id == AUTRE_VERSION


async def test_la_version_de_la_preparation_est_demandee_en_premier():
    """L'ordre est une préférence, et le service la transmet — c'est la version dans laquelle
    il travaille qui doit être nommée quand deux la reconnaissent."""
    prep = _preparation()
    versets = _Versets()
    service = _service(prep, elements=_PLAN, versets=versets)

    await service.soumettre(
        actor_account_id=AUTEUR,
        study_id=prep.id,
        diapositives=[DiapositiveSoumise("Le texte", "Hébreux 13:1", HB_13_1)],
    )

    assert versets.appels == [(58, 13, 1, None)]


# ============================================================ 4. une référence illisible


async def test_une_reference_hors_bornes_ne_bloque_que_sa_diapositive():
    """`Hb 2v29` — une vraie faute des notes du Pasteur X. Le motif dit ce qui manque **au
    corpus**, et les autres diapositives sont jugées quand même : on rend le dossier entier,
    jamais la première faute."""
    prep = _preparation()
    service = _service(prep, elements=_PLAN)

    dto = await service.soumettre(
        actor_account_id=AUTEUR,
        study_id=prep.id,
        diapositives=[
            DiapositiveSoumise("Faux", "Hb 2v29", "un texte"),
            DiapositiveSoumise("Vrai", "Hébreux 13:1", HB_13_1),
        ],
    )

    assert [c.verdict for c in dto.controles] == [ALTERE, EXACT]
    assert "18 verset" in dto.controles[0].rationale
    assert dto.record.validation == "rejete"


# ============================================================ 5. la trace


async def test_le_livrable_porte_l_empreinte_de_ce_qui_a_ete_imprime():
    """*Une décision ne vaut que sur l'objet qu'elle a regardé.* Deux documents de la même
    préparation, à deux semaines d'écart, ne sont pas le même document."""
    prep = _preparation()
    service = _service(prep, elements=_PLAN)

    premier = await service.soumettre(
        actor_account_id=AUTEUR,
        study_id=prep.id,
        diapositives=[DiapositiveSoumise("Le texte", "Hébreux 13:1", HB_13_1)],
    )
    second = await service.soumettre(
        actor_account_id=AUTEUR,
        study_id=prep.id,
        diapositives=[
            DiapositiveSoumise("Le texte", "Hébreux 13:1", "Persévérez dans l'amour")
        ],
    )

    assert premier.record.content_fingerprint != second.record.content_fingerprint
    assert premier.record.corpus_snapshot == "essai"


async def test_relire_rend_le_dossier_et_garde_la_meme_garde():
    prep = _preparation()
    depot = _Livrables()
    service = _service(prep, elements=_PLAN, livrables=depot)

    ecrit = await service.soumettre(
        actor_account_id=AUTEUR,
        study_id=prep.id,
        diapositives=[DiapositiveSoumise("Le texte", "Hébreux 13:1", HB_13_1)],
    )
    relu = await service.relire(
        actor_account_id=AUTEUR, deliverable_id=ecrit.record.id
    )

    assert relu.record.id == ecrit.record.id
    assert len(relu.controles) == 1


# ============================================================ 6. les octets


def _texte_du_pptx(octets: bytes) -> str:
    from io import BytesIO

    from pptx import Presentation

    deck = Presentation(BytesIO(octets))
    return "\n".join(
        forme.text_frame.text
        for page in deck.slides
        for forme in page.shapes
        if forme.has_text_frame
    )


def _texte_du_docx(octets: bytes) -> str:
    from io import BytesIO

    from docx import Document

    document = Document(BytesIO(octets))
    corps = "\n".join(p.text for p in document.paragraphs)
    pieds = "\n".join(
        p.text for section in document.sections for p in section.footer.paragraphs
    )
    return corps + "\n" + pieds


async def test_le_deck_produit_porte_le_texte_juge_et_rien_d_autre():
    """**La frontière, vérifiée sur le fichier lui-même** — pas sur le gabarit.

    Le type `Deck` ne peut pas porter une mise en garde ; ce test le prouve à l'autre bout de
    la chaîne, dans les octets qu'un vidéoprojecteur affichera."""
    prep = _preparation()
    service = _service(prep, elements=_PLAN)
    dossier = await service.soumettre(
        actor_account_id=AUTEUR,
        study_id=prep.id,
        diapositives=[DiapositiveSoumise("Le texte", "Hébreux 13:1", HB_13_1)],
    )

    format_, octets = await service.rendre(
        actor_account_id=AUTEUR, deliverable_id=dossier.record.id
    )

    assert format_ == "pptx"
    assert octets[:2] == b"PK"  # un .pptx est un ZIP
    contenu = _texte_du_pptx(octets)
    assert HB_13_1 in contenu
    assert "Hébreux 13:1" in contenu
    # Rien du dossier de préparation ne monte à l'écran.
    for interdit in ("proof-texting", "ia-mistral", "ne dit pas", "résiste"):
        assert interdit not in contenu


async def test_un_livrable_rejete_ne_rend_aucun_octet():
    """**La dernière porte du verrou.** Le dossier de validation revient en 201 avec ses
    verdicts — c'est ce qu'on veut montrer. Réclamer les octets de ce qui a été rejeté est
    autre chose : c'est demander ce que le contrôle existe pour ne pas produire."""
    prep = _preparation()
    service = _service(prep, elements=_PLAN)
    dossier = await service.soumettre(
        actor_account_id=AUTEUR,
        study_id=prep.id,
        diapositives=[DiapositiveSoumise("Le texte", "Hébreux 13:1", HB_13_1_ALTERE)],
    )

    with pytest.raises(LivrableNonValideError):
        await service.rendre(
            actor_account_id=AUTEUR, deliverable_id=dossier.record.id
        )


async def test_la_note_porte_ce_que_l_ecran_refuse():
    """Le couple du test précédent, de l'autre côté de la frontière : ce que le deck ne peut
    pas montrer, la note l'imprime — sinon la frontière ne protégerait rien, elle supprimerait.

    Et la mention de destination est **en pied de page**, donc sur chaque page : une page de
    garde ne survivrait ni à une capture d'écran ni à une impression recto."""
    prep = _preparation()
    etude = _Etude(_dossier_d_etude(prep))
    service = _service(prep, elements=_PLAN, etude=etude)
    dossier = await service.soumettre(
        actor_account_id=AUTEUR, study_id=prep.id, kind="note"
    )

    format_, octets = await service.rendre(
        actor_account_id=AUTEUR, deliverable_id=dossier.record.id
    )

    assert format_ == "docx"
    contenu = _texte_du_docx(octets)
    assert "Le texte ne dit pas quand l'Esprit est donné." in contenu
    assert "ia-mistral" in contenu  # la signature de la curation remonte jusqu'au papier
    assert "s'adressent au prédicateur" in contenu
    assert etude.lectures == [prep.id]  # le dossier a été rejoué, pas reconstruit


async def test_le_service_du_livrable_n_a_aucun_port_de_reservation():
    """**Générer ne peut rien consommer**, et ce n'est pas une intention : le service n'a aucun
    moyen de compter quoi que ce soit. Un port absent ne s'appelle pas par accident."""
    from dataclasses import fields

    champs = {f.name for f in fields(UrimDeliverableService)}
    assert "reservations" not in champs
    assert "resolver" not in champs  # ni modèle : un document n'appelle personne


async def test_l_horloge_du_service_date_le_livrable():
    prep = _preparation()
    service = _service(prep, elements=_PLAN)
    dto = await service.soumettre(actor_account_id=AUTEUR, study_id=prep.id, kind="note")
    assert dto.record.generated_at == datetime(2026, 8, 10, tzinfo=UTC)
