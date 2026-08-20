"""Les deux documents portent la charte — **et ils sortent du produit**.

Le `.pptx` monte au mur d'une assemblée, la note part au bureau du pasteur et souvent chez
l'imprimeur. Ce sont les seuls objets d'Urim que des gens qui n'ont jamais ouvert
l'application vont voir.

Ce que ces tests gardent n'est pas le goût — c'est qu'aucun des deux ne reparte avec l'identité
de son outil : le bleu Office de `python-pptx`, le Calibri de Word. Un document qui garde ces
traits n'est celui de personne.
"""

from __future__ import annotations

from io import BytesIO

import pytest

from app.contexts.urim.deliverable.domain.documents import Deck, Diapositive, Note
from app.contexts.urim.deliverable.infrastructure import charte
from app.contexts.urim.deliverable.infrastructure.renderers import rendre_deck, rendre_note

pptx = pytest.importorskip("pptx", reason="le rendu dépend de python-pptx")
docx = pytest.importorskip("docx", reason="le rendu dépend de python-docx")


def _deck() -> bytes:
    return rendre_deck(
        Deck(
            titre="La prééminence du Christ",
            diapositives=(
                Diapositive(
                    titre="",
                    reference="Colossiens 1:15",
                    texte_projete="Il est l'image du Dieu invisible.",
                ),
                Diapositive(
                    titre="Tout subsiste en lui",
                    reference="Colossiens 1:17",
                    texte_projete="Il est avant toutes choses.",
                ),
            ),
        )
    )


def _presentation():
    from pptx import Presentation

    return Presentation(BytesIO(_deck()))


def test_la_couverture_est_marine_et_les_textes_sur_sable():
    """Le fond ne vient pas du gabarit Office : il est posé, diapositive par diapositive."""
    pages = list(_presentation().slides)

    assert str(pages[0].background.fill.fore_color.rgb) == "003049"
    assert all(str(p.background.fill.fore_color.rgb) == "F7F4E4" for p in pages[1:])


def test_chaque_diapositive_dit_d_ou_elle_vient():
    """La signature en pied, sur la couverture comme sur les textes.

    Discrète et jamais un logo : une image embarquée triple le poids du fichier et sort floue
    au vidéoprojecteur."""
    pages = list(_presentation().slides)
    textes = [
        " ".join(f.text_frame.text for f in page.shapes if f.has_text_frame)
        for page in pages
    ]

    assert charte.SIGNATURE in textes[0]
    assert any(charte.SIGNATURE in t for t in textes[1:])


def test_une_diapositive_titree_rappelle_sa_reference_en_pied():
    """Une projection sans référence est une citation invérifiable depuis le banc.

    ⚠️ Et **une seule fois** : quand la référence est déjà l'intitulé, la répéter en pied se
    voit du fond de la salle."""
    pages = list(_presentation().slides)
    sans_titre = " ".join(
        f.text_frame.text for f in pages[1].shapes if f.has_text_frame
    )
    avec_titre = " ".join(
        f.text_frame.text for f in pages[2].shapes if f.has_text_frame
    )

    # La comparaison est **insensible à la casse** : sur une diapositive sans intitulé, la
    # référence *est* l'intitulé, et un intitulé se pose en capitales.
    assert sans_titre.lower().count("colossiens 1:15") == 1
    assert avec_titre.lower().count("colossiens 1:17") == 1
    assert "TOUT SUBSISTE EN LUI" in avec_titre


def test_le_texte_projete_se_reduit_quand_il_s_allonge():
    """Un corps fixe fait déborder la diapositive, et un débordement se voit avant le texte."""
    assert charte.corps_du_verset("Un verset court.") == 30
    assert charte.corps_du_verset("m" * 300) == 24
    assert charte.corps_du_verset("m" * 600) == 20


def _note():
    from docx import Document

    return Document(
        BytesIO(
            rendre_note(
                Note(
                    titre="La prééminence du Christ",
                    reference="Colossiens 1:15-20",
                    unite="Christologie : la prééminence du Christ",
                    motif_unite="Le passage développe la prééminence.",
                    plan=(),
                    versets=(),
                    version="LSG 1910",
                    pesees=(),
                    axe_retenu="christologie",
                    mises_en_garde=(),
                    faisabilites=(),
                    resistances=(),
                    appuis=(),
                    original=(),
                    ecartees=(),
                    signature=None,
                    corpus_snapshot="0631af9c",
                )
            )
        )
    )


def test_la_note_ne_repart_pas_avec_les_styles_de_word():
    document = _note()

    assert document.styles["Normal"].font.name == charte.SERIF
    assert str(document.styles["Heading 1"].font.color.rgb) == "003049"
    assert str(document.styles["Heading 2"].font.color.rgb) == "CC3C1F"


def test_le_pied_porte_la_mention_et_la_signature():
    """La mention de destination est en pied de **chaque** page — une page de garde ne survit
    ni à une capture d'écran, ni à un partage partiel, ni à une impression recto."""
    pied = _note().sections[0].footer.paragraphs[0].text

    assert "les mises en garde s'adressent au prédicateur" in pied
    assert charte.SIGNATURE in pied


# ------------------------------------------------------------------ le titre du document


class _Trace(tuple):
    pass


def _etude(theme: str | None, axe: str = "christologie"):
    """Le strict nécessaire pour titrer — le reste de la note ne s'en sert pas."""
    from types import SimpleNamespace

    return SimpleNamespace(
        record=SimpleNamespace(
            theme=theme,
            axis_code=axe,
            plan_source="thematique",
            subject_matter="doctrinal",
        ),
        pericope_label="Adresse et action de grâces pour les Colossiens",
        resolved_label="Colossiens 1:1-14",
    )


def test_le_gabarit_du_moteur_ne_titre_pas_le_document():
    """🔴 La note s'intitulait « christologie, en thematique doctrinal ».

    C'est un code d'axe et deux codes de forme recollés. Imprimé en tête d'une fiche qu'on
    emporte en chaire, il fait passer un état interne pour une intention — et le moteur dit
    lui-même *« un thème, jamais un titre »*."""
    from app.contexts.urim.deliverable.application.service import _titre_de

    assert (
        _titre_de(_etude("christologie, en thematique doctrinal"))
        == "Adresse et action de grâces pour les Colossiens"
    )


def test_une_phrase_du_pasteur_titre_le_document():
    """Dès qu'il a réécrit, c'est sa voix — et elle passe avant l'unité relue."""
    from app.contexts.urim.deliverable.application.service import _titre_de

    assert _titre_de(_etude("Jésus, notre vrai trésor")) == "Jésus, notre vrai trésor"


def test_sans_theme_ni_unite_il_reste_la_reference():
    from app.contexts.urim.deliverable.application.service import _titre_de

    etude = _etude(None)
    etude.pericope_label = None

    assert _titre_de(etude) == "Colossiens 1:1-14"
