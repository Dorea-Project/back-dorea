"""DOREA-024 au dépôt de sermon — le `kind` est déclaré, les octets sont vérifiés.

`POST /api/mobile/sermons/tenants/{tid}/upload` reçoit le format en **paramètre de requête** :
le client dit ce qu'il veut. Sans recoupement, des octets arbitraires déclarés `pdf` entraient
dans `pypdf`, et déclarés `pptx` dans `zipfile` — un parseur nourri d'entrée non fiable.

Chaque cas ci-dessous présente le **couple** : le fichier légitime, que le dépôt accepte, et sa
jumelle fautive, qu'il refuse. C'est le couple qui prouve quelque chose — une garde qui rejette
tout est aussi inutile qu'une garde absente, et elle se remarque plus tard, quand un pasteur voit
son vrai PDF refusé (même raisonnement qu'en tête de `tests/contexts/urim/test_schema_urim.py`).

Le couple le plus parlant est celui du PPTX déclaré `pdf` : **les mêmes octets** sont le fichier
légitime d'un cas et la fraude de l'autre. Ce n'est donc pas la qualité du fichier qui est jugée,
c'est l'accord entre ce qu'il est et ce qu'il prétend être.
"""

from datetime import UTC, date, datetime
from io import BytesIO
from uuid import uuid4

import pytest

from app.contexts.groups.application.group_access import GroupAccessPolicy
from app.contexts.iam.application.ports import OwnershipChecker
from app.contexts.iam.domain.aggregates import Membership
from app.contexts.iam.domain.entities import RoleAssignment
from app.contexts.iam.domain.enums import MembershipStatus, RoleCode
from app.contexts.iam.domain.repositories import MembershipRepository
from app.contexts.sermon.application.commands.deposit import DepositSermon
from app.contexts.sermon.application.ports import SermonTextExtractor
from app.contexts.sermon.application.upload_guard import ensure_declared_kind
from app.contexts.sermon.domain.enums import SermonSourceKind
from app.contexts.sermon.domain.errors import (
    SermonFileTypeMismatchError,
    UnsupportedSermonFormatError,
)
from app.contexts.sermon.domain.repositories import SermonRepository
from app.contexts.sermon.infrastructure.extractor import CompositeTextExtractor

_NOW = datetime(2026, 1, 4, tzinfo=UTC)
_SUNDAY = date(2026, 1, 4)


# --- Les fichiers : deux vrais, et de quoi mentir ---


def _real_pdf() -> bytes:
    from pypdf import PdfWriter

    buffer = BytesIO()
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.write(buffer)
    return buffer.getvalue()


def _real_pptx() -> bytes:
    from pptx import Presentation

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "La grâce du Père"
    buffer = BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


# --- La garde elle-même, format par format ---


def test_un_vrai_pdf_est_accepte():
    """Le jumeau légitime — une garde qui refuse le bon fichier est une régression."""
    ensure_declared_kind(_real_pdf(), SermonSourceKind.PDF)  # pas d'erreur


def test_un_pptx_declare_pdf_est_refuse():
    """Les mêmes octets, légitimes en `pptx`, mentent en `pdf` : c'est l'accord qui est jugé."""
    with pytest.raises(SermonFileTypeMismatchError):
        ensure_declared_kind(_real_pptx(), SermonSourceKind.PDF)


def test_un_vrai_pptx_est_accepte():
    ensure_declared_kind(_real_pptx(), SermonSourceKind.PPTX)  # pas d'erreur


def test_des_octets_arbitraires_declares_pptx_sont_refuses():
    """Avant, ceux-ci descendaient jusqu'à `zipfile` ; ils s'arrêtent maintenant à la porte."""
    with pytest.raises(SermonFileTypeMismatchError):
        ensure_declared_kind(b"ceci n'est pas un pptx", SermonSourceKind.PPTX)


def test_un_pdf_declare_pptx_est_refuse():
    with pytest.raises(SermonFileTypeMismatchError):
        ensure_declared_kind(_real_pdf(), SermonSourceKind.PPTX)


def test_un_fichier_vide_ne_ressemble_a_aucun_format_signe():
    for kind in (SermonSourceKind.PDF, SermonSourceKind.PPTX):
        with pytest.raises(SermonFileTypeMismatchError):
            ensure_declared_kind(b"", kind)


def test_le_texte_n_a_pas_de_signature_et_passe_tel_quel():
    """Le couple inverse : la garde ne parle que des formats qu'elle sait reconnaître.

    Du texte n'a pas d'en-tête et n'a pas de parseur derrière — rien à protéger, rien à exiger."""
    ensure_declared_kind(b"Aucune condamnation", SermonSourceKind.TEXT)
    ensure_declared_kind(b"\xff\xfe des octets qui ne se decodent pas", SermonSourceKind.TEXT)


def test_l_erreur_de_signature_est_un_415_qui_nomme_le_format():
    """Le refus doit être lisible côté client — sinon il se traite comme une panne."""
    with pytest.raises(SermonFileTypeMismatchError) as raised:
        ensure_declared_kind(b"MZ\x90\x00", SermonSourceKind.PDF)
    assert raised.value.http_status == 415
    assert raised.value.details == {"kind": "pdf"}


def test_l_audio_reste_un_format_non_pris_en_charge():
    """S-6 n'est pas branché — et l'absence de table est ce qui le dit (fail-closed).

    Ce n'est pas un mensonge sur les octets : c'est un format qu'on n'accepte pas encore. Les
    deux refus ne se confondent donc pas, et le jour où l'audio arrivera, il faudra d'abord
    écrire comment on le reconnaît (D8 : la durée lue, jamais déclarée)."""
    with pytest.raises(UnsupportedSermonFormatError) as raised:
        ensure_declared_kind(b"\x00\x00\x00\x20ftypM4A ", SermonSourceKind.AUDIO)
    # Et pas un mensonge de signature : les deux refus restent distincts pour le client.
    assert not isinstance(raised.value, SermonFileTypeMismatchError)


# --- Le dépôt complet : le parseur ne voit que ce qui a passé la porte ---


class _FakeOwnership(OwnershipChecker):
    async def is_active_owner(self, account_id, tenant_id):
        return False


class _FakeMemberships(MembershipRepository):
    def __init__(self, items=()):
        self._m = list(items)

    async def get_active(self, account_id, tenant_id):
        return next(
            (
                m for m in self._m
                if m.account_id == account_id and m.tenant_id == tenant_id and not m.is_closed
            ),
            None,
        )

    async def list_active_by_account(self, account_id):
        return [m for m in self._m if m.account_id == account_id]

    async def count_active_group_leaders(self, tenant_id, group_id):
        return 0


class _FakeSermons(SermonRepository):
    def __init__(self):
        self.added = []

    async def add(self, sermon):
        self.added.append(sermon)

    async def get(self, sermon_id):
        return None

    async def save(self, sermon):
        pass

    async def list_by_tenant(self, tenant_id):
        return []


class _SpyExtractor(SermonTextExtractor):
    """Tient le compte des octets qui lui parviennent — c'est ce compte qui prouve la garde."""

    def __init__(self):
        self.calls = []

    async def extract(self, data, *, kind):
        self.calls.append((data, kind))
        return "Le fils prodigue revint vers son père."


def _pastor_command(extractor):
    tenant, pastor = uuid4(), uuid4()
    membership = Membership(
        id=uuid4(), account_id=pastor, tenant_id=tenant,
        status=MembershipStatus.CONFIRMED_MEMBER, last_transition_at=_NOW,
        role_assignments=[
            RoleAssignment(
                id=uuid4(), role=RoleCode.PASTOR, group_id=None,
                assigned_at=_NOW, assigned_by_account_id=uuid4(),
            )
        ],
    )
    access = GroupAccessPolicy(_FakeOwnership(), _FakeMemberships([membership]))
    sermons = _FakeSermons()
    command = DepositSermon(sermons, access, None, extractor, clock=lambda: _NOW)
    return command, sermons, tenant, pastor


async def _deposit(command, tenant, pastor, data, kind):
    return await command.execute_file(
        actor_account_id=pastor, tenant_id=tenant, title="Depuis un fichier",
        data=data, kind=kind, preached_on=_SUNDAY,
    )


async def test_le_depot_d_un_vrai_pdf_atteint_l_extracteur():
    """Le jumeau légitime, au niveau du dépôt : le chemin nominal reste ouvert."""
    extractor = _SpyExtractor()
    command, sermons, tenant, pastor = _pastor_command(extractor)

    dto = await _deposit(command, tenant, pastor, _real_pdf(), SermonSourceKind.PDF)

    assert dto.source_kind == "pdf" and dto.raw_text.startswith("Le fils prodigue")
    assert extractor.calls[0][1] is SermonSourceKind.PDF
    assert len(sermons.added) == 1


async def test_le_depot_menteur_est_refuse_avant_que_le_parseur_ne_voie_rien():
    """Le point de tout l'exercice : le refus tombe **avant** l'extraction, pas pendant.

    Que `pypdf` échoue de lui-même sur du n'importe quoi est un accident heureux, pas une
    garantie — et le jour où l'extracteur appelle un service facturé (S-6, D8bis), l'échec
    du parseur arrive une facture trop tard."""
    extractor = _SpyExtractor()
    command, sermons, tenant, pastor = _pastor_command(extractor)

    with pytest.raises(SermonFileTypeMismatchError):
        await _deposit(command, tenant, pastor, b"<?php echo 1; ?>" * 4, SermonSourceKind.PDF)

    assert extractor.calls == []  # le parseur n'a rien reçu
    assert sermons.added == []  # et rien n'est entré au dépôt


async def test_un_depot_audio_ne_touche_pas_davantage_l_extracteur():
    """S-6 non branché : le format est refusé au dépôt, l'extracteur n'est pas sollicité."""
    extractor = _SpyExtractor()
    command, sermons, tenant, pastor = _pastor_command(extractor)

    with pytest.raises(UnsupportedSermonFormatError):
        await _deposit(command, tenant, pastor, b"\x00\x00\x00\x20ftypM4A ", SermonSourceKind.AUDIO)

    assert extractor.calls == [] and sermons.added == []


async def test_l_extracteur_reel_refuse_toujours_l_audio_de_son_cote():
    """La garde ne remplace pas l'adaptateur : les deux refusent, indépendamment."""
    with pytest.raises(UnsupportedSermonFormatError):
        await CompositeTextExtractor().extract(b"...", kind=SermonSourceKind.AUDIO)
