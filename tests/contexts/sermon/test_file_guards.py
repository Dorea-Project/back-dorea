"""DOREA-011 — un fichier déposé ne doit pas coûter plus cher qu'il ne pèse.

Le plafond d'upload (15 Mo) borne ce qui **entre**. Il ne borne pas ce que ça **coûte** :
un `.pptx` est un ZIP, et quinze mégaoctets compressés peuvent en peser des milliers une
fois décompressés. Un PDF peut porter des dizaines de milliers de pages.

Ces tests construisent de vraies bombes — petites, mais avec le même taux d'expansion que
les grandes — et vérifient que l'extracteur refuse **avant** de dépouiller.
"""

import zipfile
from io import BytesIO

import pytest

from app.contexts.sermon.domain.enums import SermonSourceKind
from app.contexts.sermon.domain.errors import SermonFileTooComplexError
from app.contexts.sermon.infrastructure.extractor import CompositeTextExtractor

_extractor = CompositeTextExtractor()


def _zip_bomb(*, members: int = 40, size: int = 2 * 1024 * 1024) -> bytes:
    """Un ZIP minuscule qui se déploie en dizaines de Mio — des zéros compressent trop bien."""
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for index in range(members):
            archive.writestr(f"ppt/slides/slide{index}.xml", b"\0" * size)
    return buffer.getvalue()


async def test_une_bombe_zip_est_refusee_avant_ouverture():
    bomb = _zip_bomb()
    # La bombe est bien minuscule *compressée* — le plafond d'upload ne la voit pas passer.
    assert len(bomb) < 1024 * 1024

    with pytest.raises(SermonFileTooComplexError):
        await _extractor.extract(bomb, kind=SermonSourceKind.PPTX)


async def test_le_taux_d_expansion_est_le_critere():
    """Ce n'est pas la taille seule : c'est le rapport entre ce qui entre et ce qui sort."""
    with pytest.raises(SermonFileTooComplexError) as raised:
        await _extractor.extract(_zip_bomb(members=10), kind=SermonSourceKind.PPTX)
    details = raised.value.details
    assert details["decompresse_octets"] > details["compresse"] * 100


async def test_un_pptx_illisible_est_un_format_invalide_pas_une_bombe():
    """Un fichier qui n'est pas un ZIP échoue proprement, sans confusion de diagnostic."""
    from app.contexts.sermon.domain.errors import UnsupportedSermonFormatError

    with pytest.raises(UnsupportedSermonFormatError):
        await _extractor.extract(b"ceci n'est pas un pptx", kind=SermonSourceKind.PPTX)


async def test_un_pptx_honnete_passe():
    """Le jumeau légitime : un vrai deck, modeste, n'est pas gêné par la garde."""
    from pptx import Presentation

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "Il n'y a donc maintenant aucune condamnation"
    buffer = BytesIO()
    deck.save(buffer)

    text = await _extractor.extract(buffer.getvalue(), kind=SermonSourceKind.PPTX)
    assert "aucune condamnation" in text


async def test_le_texte_rendu_est_borne():
    """DOREA-026 — ce qui part vers le modèle a une borne (ici sur l'entrée texte)."""
    enorme = ("a" * 1000 + "\n") * 2000  # ~2 Mo de texte
    rendu = await _extractor.extract(enorme.encode("utf-8"), kind=SermonSourceKind.TEXT)
    assert len(rendu) <= 400_000
