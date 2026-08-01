"""Module Sermon — S-0 : le socle (dépôt texte + cycle de vie brouillon → approuvé → publié).

Le pasteur dépose son sermon (autorité `PUBLISH_SERMON`) ; il l'approuve (son onction) puis le
publie. Rien de non approuvé n'est publiable. Un membre ordinaire ne peut pas déposer.
"""

from datetime import UTC, date, datetime
from io import BytesIO
from uuid import uuid4

import pytest

from app.contexts.groups.application.group_access import GroupAccessPolicy
from app.contexts.groups.domain.errors import UnauthorizedGroupActionError
from app.contexts.iam.application.ports import OwnershipChecker
from app.contexts.iam.domain.aggregates import Membership
from app.contexts.iam.domain.entities import RoleAssignment
from app.contexts.iam.domain.enums import MembershipStatus, RoleCode
from app.contexts.iam.domain.repositories import MembershipRepository
from app.contexts.sermon.application.commands.companion import (
    AdvanceCompanion,
    AnswerAttendance,
    StartCompanion,
)
from app.contexts.sermon.application.commands.deposit import DepositSermon
from app.contexts.sermon.application.commands.manage import ApproveSermon, PublishSermon
from app.contexts.sermon.application.ports import (
    CapsuleFeedPort,
    CulteAttendancePort,
    SermonDigester,
    SermonTextExtractor,
)
from app.contexts.sermon.application.queries.list_sermons import GetSermon, ListTenantSermons
from app.contexts.sermon.domain.aggregates import Sermon
from app.contexts.sermon.domain.digest import Capsule, CompanionQuestion, SermonDigest
from app.contexts.sermon.domain.enums import SermonSourceKind, SermonStatus
from app.contexts.sermon.domain.errors import (
    NotAChurchMemberError,
    NotSessionOwnerError,
    SermonContentRequiredError,
    SermonNotEditableError,
    SermonNotFoundError,
    SermonNotPublishedError,
    SermonTitleRequiredError,
    UnsupportedSermonFormatError,
)
from app.contexts.sermon.domain.repositories import (
    CompanionSessionRepository,
    SermonRepository,
)
from app.contexts.sermon.infrastructure.digester import KeywordSermonDigester
from app.contexts.sermon.infrastructure.extractor import CompositeTextExtractor

_NOW = datetime(2026, 1, 4, tzinfo=UTC)  # un dimanche
_SUNDAY = date(2026, 1, 4)


# --- fakes ---


class _FakeOwnership(OwnershipChecker):
    def __init__(self, owners=()):
        self._owners = set(owners)

    async def is_active_owner(self, account_id, tenant_id):
        return (account_id, tenant_id) in self._owners


class _FakeMemberships(MembershipRepository):
    def __init__(self, items=()):
        self._m = list(items)

    async def get_active(self, account_id, tenant_id):
        return next(
            (
                m for m in self._m
                if m.account_id == account_id and m.tenant_id == tenant_id and not m.is_closed
            ),
            None,
        )

    async def list_active_by_account(self, account_id):
        return [m for m in self._m if m.account_id == account_id]

    async def count_active_group_leaders(self, tenant_id, group_id):
        return 0


class _FakeSermons(SermonRepository):
    def __init__(self, items=()):
        self._s = list(items)

    async def add(self, sermon):
        self._s.append(sermon)

    async def get(self, sermon_id):
        return next((s for s in self._s if s.id == sermon_id), None)

    async def save(self, sermon):
        pass  # muté en mémoire (même instance)

    async def list_by_tenant(self, tenant_id):
        return [s for s in self._s if s.tenant_id == tenant_id]


def _member(account, tenant, *roles) -> Membership:
    ras = [
        RoleAssignment(
            id=uuid4(), role=r, group_id=g, assigned_at=_NOW, assigned_by_account_id=uuid4()
        )
        for (r, g) in roles
    ]
    return Membership(
        id=uuid4(), account_id=account, tenant_id=tenant,
        status=MembershipStatus.CONFIRMED_MEMBER, last_transition_at=_NOW, role_assignments=ras,
    )


def _access(memberships, *, owners=()):
    return GroupAccessPolicy(_FakeOwnership(owners), memberships)


def _pastor(tenant):
    pastor = uuid4()
    ms = _FakeMemberships([_member(pastor, tenant, (RoleCode.PASTOR, None))])
    return pastor, ms


def _draft(tenant, author) -> Sermon:
    return Sermon.deposit(
        id=uuid4(), tenant_id=tenant, author_account_id=author,
        title="L'amour qui pardonne", raw_text="Le père courut vers son fils…",
        preached_on=_SUNDAY, now=_NOW, reference="Luc 15.11-32",
    )


# --- Le domaine : dépôt + cycle de vie ---


def test_deposit_requires_a_title():
    with pytest.raises(SermonTitleRequiredError):
        Sermon.deposit(
            id=uuid4(), tenant_id=uuid4(), author_account_id=uuid4(),
            title="   ", raw_text="du texte", preached_on=_SUNDAY, now=_NOW,
        )


def test_deposit_requires_content():
    with pytest.raises(SermonContentRequiredError):
        Sermon.deposit(
            id=uuid4(), tenant_id=uuid4(), author_account_id=uuid4(),
            title="Un titre", raw_text="   ", preached_on=_SUNDAY, now=_NOW,
        )


def test_a_sermon_is_born_a_draft():
    s = _draft(uuid4(), uuid4())
    assert s.status is SermonStatus.DRAFT and s.approved_at is None
    assert s.reference == "Luc 15.11-32"


def test_approve_then_publish_walks_the_lifecycle():
    s = _draft(uuid4(), uuid4())
    s.approve(now=_NOW)
    assert s.status is SermonStatus.APPROVED and s.approved_at == _NOW
    s.publish(now=_NOW)
    assert s.status is SermonStatus.PUBLISHED and s.is_published


def test_cannot_publish_before_approval():
    s = _draft(uuid4(), uuid4())
    with pytest.raises(SermonNotEditableError):
        s.publish(now=_NOW)


def test_cannot_approve_twice():
    s = _draft(uuid4(), uuid4())
    s.approve(now=_NOW)
    with pytest.raises(SermonNotEditableError):
        s.approve(now=_NOW)


# --- Le pasteur (mobile) ---


async def test_pastor_deposits_a_sermon():
    tenant = uuid4()
    pastor, ms = _pastor(tenant)
    sermons = _FakeSermons()
    cmd = DepositSermon(sermons, _access(ms), clock=lambda: _NOW)
    dto = await cmd.execute(
        actor_account_id=pastor, tenant_id=tenant,
        title="La grâce", content="Par grâce vous êtes sauvés…", preached_on=_SUNDAY,
    )
    assert dto.status == "draft" and dto.author_account_id == pastor
    assert len(sermons._s) == 1


async def test_an_ordinary_member_cannot_deposit_a_sermon():
    tenant, stranger = uuid4(), uuid4()
    ms = _FakeMemberships([_member(stranger, tenant)])  # aucun rôle
    cmd = DepositSermon(_FakeSermons(), _access(ms), clock=lambda: _NOW)
    with pytest.raises(UnauthorizedGroupActionError):
        await cmd.execute(
            actor_account_id=stranger, tenant_id=tenant,
            title="x", content="y", preached_on=_SUNDAY,
        )


async def test_owner_can_deposit_even_without_a_role():
    tenant, owner = uuid4(), uuid4()
    ms = _FakeMemberships([_member(owner, tenant)])
    cmd = DepositSermon(_FakeSermons(), _access(ms, owners=[(owner, tenant)]), clock=lambda: _NOW)
    dto = await cmd.execute(
        actor_account_id=owner, tenant_id=tenant,
        title="Le bon berger", content="Je suis le bon berger…", preached_on=_SUNDAY,
    )
    assert dto.status == "draft"


async def test_pastor_approves_then_publishes():
    tenant = uuid4()
    pastor, ms = _pastor(tenant)
    s = _draft(tenant, pastor)
    sermons = _FakeSermons([s])
    await ApproveSermon(sermons, _access(ms), clock=lambda: _NOW).execute(
        actor_account_id=pastor, sermon_id=s.id
    )
    dto = await PublishSermon(sermons, _access(ms), clock=lambda: _NOW).execute(
        actor_account_id=pastor, sermon_id=s.id
    )
    assert dto.status == "published"


async def test_publishing_an_unknown_sermon_is_404():
    tenant = uuid4()
    pastor, ms = _pastor(tenant)
    cmd = PublishSermon(_FakeSermons(), _access(ms), clock=lambda: _NOW)
    with pytest.raises(SermonNotFoundError):
        await cmd.execute(actor_account_id=pastor, sermon_id=uuid4())


async def test_list_shows_the_church_sermons_to_the_pastor():
    tenant = uuid4()
    pastor, ms = _pastor(tenant)
    sermons = _FakeSermons([_draft(tenant, pastor), _draft(tenant, pastor)])
    dtos = await ListTenantSermons(sermons, _access(ms)).execute(
        actor_account_id=pastor, tenant_id=tenant
    )
    assert len(dtos) == 2


async def test_a_non_keeper_cannot_list_sermons():
    tenant, stranger = uuid4(), uuid4()
    ms = _FakeMemberships([_member(stranger, tenant)])
    with pytest.raises(UnauthorizedGroupActionError):
        await ListTenantSermons(_FakeSermons(), _access(ms)).execute(
            actor_account_id=stranger, tenant_id=tenant
        )


async def test_get_sermon_returns_it_to_the_keeper():
    tenant = uuid4()
    pastor, ms = _pastor(tenant)
    s = _draft(tenant, pastor)
    dto = await GetSermon(_FakeSermons([s]), _access(ms)).execute(
        actor_account_id=pastor, sermon_id=s.id
    )
    assert dto.id == s.id and dto.raw_text == s.raw_text


# --- S-1 : la digestion IA (un appel au dépôt, gelé à l'approbation) ---


class _FakeDigester(SermonDigester):
    def __init__(self):
        self.calls = []

    async def digest(self, text, *, title, reference):
        self.calls.append((text, title, reference))
        return SermonDigest(
            summary="Dieu court vers nous.",
            key_points=("Le père guette", "Il court", "Il restaure"),
            capsules=(Capsule(title="Le retour", body="Le fils revint…"),),
            questions=(CompanionQuestion(prompt="Où en es-tu ?", guidance="Dieu t'attend."),),
        )


async def test_deposit_generates_a_digest_in_one_call():
    tenant = uuid4()
    pastor, ms = _pastor(tenant)
    sermons = _FakeSermons()
    digester = _FakeDigester()
    cmd = DepositSermon(sermons, _access(ms), digester, clock=lambda: _NOW)
    dto = await cmd.execute(
        actor_account_id=pastor, tenant_id=tenant,
        title="Le fils prodigue", content="Le père courut vers son fils…", preached_on=_SUNDAY,
    )
    assert len(digester.calls) == 1  # un seul appel IA au dépôt
    assert dto.digest is not None and dto.digest.summary == "Dieu court vers nous."
    assert sermons._s[0].digest is not None  # attaché à l'agrégat


async def test_deposit_without_a_digester_leaves_no_digest():
    tenant = uuid4()
    pastor, ms = _pastor(tenant)
    cmd = DepositSermon(_FakeSermons(), _access(ms), clock=lambda: _NOW)
    dto = await cmd.execute(
        actor_account_id=pastor, tenant_id=tenant,
        title="Sobre", content="Un texte.", preached_on=_SUNDAY,
    )
    assert dto.digest is None


async def test_keyword_digester_is_deterministic_without_a_key():
    digester = KeywordSermonDigester()
    text = "Dieu est amour. Il donne sa vie. Nous sommes sauvés. Il faut croire."
    d = await digester.digest(text, title="La grâce", reference=None)
    assert d.summary and d.key_points and d.capsules and d.questions
    # tiré du texte lui-même (aucune invention)
    assert d.key_points[0] == "Dieu est amour."


# --- S-3 : le compagnon (arbre déterministe, deux branches) ---


_DIGEST = SermonDigest(
    summary="Dieu court vers nous.",
    key_points=("Le père guette", "Il court", "Il restaure"),
    capsules=(Capsule(title="Le retour", body="Le fils revint…"),),
    questions=(
        CompanionQuestion(prompt="Où en es-tu ?", guidance="Dieu t'attend."),
        CompanionQuestion(prompt="Que gardes-tu ?", guidance="Sa grâce demeure."),
    ),
)


class _FakeSessions(CompanionSessionRepository):
    def __init__(self):
        self._x = []

    async def add(self, session):
        self._x.append(session)

    async def get(self, session_id):
        return next((s for s in self._x if s.id == session_id), None)

    async def save(self, session):
        pass  # muté en mémoire (même instance)

    async def find_active(self, member_account_id, sermon_id):
        return next(
            (
                s for s in self._x
                if s.member_account_id == member_account_id
                and s.sermon_id == sermon_id
                and not s.is_completed
            ),
            None,
        )


def _published(tenant, author) -> Sermon:
    s = _draft(tenant, author)
    s.attach_digest(_DIGEST)
    s.approve(now=_NOW)
    s.publish(now=_NOW)
    return s


def _member_ms(tenant):
    member = uuid4()
    ms = _FakeMemberships([_member(member, tenant)])
    return member, ms


def _start(sessions, sermons, ms):
    return StartCompanion(sessions, sermons, ms, clock=lambda: _NOW)


def _answer(sessions, sermons):
    return AnswerAttendance(sessions, sermons, clock=lambda: _NOW)


def _advance(sessions, sermons):
    return AdvanceCompanion(sessions, sermons, clock=lambda: _NOW)


async def test_start_companion_asks_the_entry_question():
    tenant = uuid4()
    member, ms = _member_ms(tenant)
    sessions, sermons = _FakeSessions(), _FakeSermons([_published(tenant, uuid4())])
    card = await _start(sessions, sermons, ms).execute(
        actor_account_id=member, sermon_id=sermons._s[0].id
    )
    assert card.stage == "attendance" and "culte" in card.prompt and not card.done
    assert len(sessions._x) == 1


async def test_start_is_refused_on_an_unpublished_sermon():
    tenant = uuid4()
    member, ms = _member_ms(tenant)
    draft = _draft(tenant, uuid4())  # brouillon
    with pytest.raises(SermonNotPublishedError):
        await _start(_FakeSessions(), _FakeSermons([draft]), ms).execute(
            actor_account_id=member, sermon_id=draft.id
        )


async def test_a_non_member_cannot_start_the_companion():
    tenant = uuid4()
    sermon = _published(tenant, uuid4())
    with pytest.raises(NotAChurchMemberError):
        await _start(_FakeSessions(), _FakeSermons([sermon]), _FakeMemberships()).execute(
            actor_account_id=uuid4(), sermon_id=sermon.id
        )


async def test_yes_walks_the_consolidation_branch_to_a_closing():
    tenant = uuid4()
    member, ms = _member_ms(tenant)
    sermon = _published(tenant, uuid4())
    sessions, sermons = _FakeSessions(), _FakeSermons([sermon])
    start = await _start(sessions, sermons, ms).execute(
        actor_account_id=member, sermon_id=sermon.id
    )
    sid = start.session_id
    c0 = await _answer(sessions, sermons).execute(
        actor_account_id=member, session_id=sid, attended=True
    )
    assert c0.stage == "consolidation" and c0.prompt == "Où en es-tu ?" and c0.guidance
    c1 = await _advance(sessions, sermons).execute(actor_account_id=member, session_id=sid)
    assert c1.stage == "consolidation" and c1.index == 1
    end = await _advance(sessions, sermons).execute(actor_account_id=member, session_id=sid)
    assert end.stage == "closing" and end.done


async def test_no_walks_the_teaching_branch_the_ministry_not_a_reproach():
    tenant = uuid4()
    member, ms = _member_ms(tenant)
    sermon = _published(tenant, uuid4())
    sessions, sermons = _FakeSessions(), _FakeSermons([sermon])
    start = await _start(sessions, sermons, ms).execute(
        actor_account_id=member, sermon_id=sermon.id
    )
    card = await _answer(sessions, sermons).execute(
        actor_account_id=member, session_id=start.session_id, attended=False
    )
    # branche enseignement : les points essentiels, pas un jugement
    assert card.stage == "teaching" and card.prompt == "Le père guette" and card.total == 3


async def test_starting_again_resumes_the_same_session():
    tenant = uuid4()
    member, ms = _member_ms(tenant)
    sermon = _published(tenant, uuid4())
    sessions, sermons = _FakeSessions(), _FakeSermons([sermon])
    start = _start(sessions, sermons, ms)
    first = await start.execute(actor_account_id=member, sermon_id=sermon.id)
    await _answer(sessions, sermons).execute(
        actor_account_id=member, session_id=first.session_id, attended=True
    )
    again = await start.execute(actor_account_id=member, sermon_id=sermon.id)
    assert again.session_id == first.session_id  # reprise, pas de doublon
    assert again.stage == "consolidation"  # on retrouve là où on en était
    assert len(sessions._x) == 1


async def test_a_companion_session_is_private():
    tenant = uuid4()
    member, ms = _member_ms(tenant)
    sermon = _published(tenant, uuid4())
    sessions, sermons = _FakeSessions(), _FakeSermons([sermon])
    start = await _start(sessions, sermons, ms).execute(
        actor_account_id=member, sermon_id=sermon.id
    )
    with pytest.raises(NotSessionOwnerError):  # un autre membre ne peut pas y toucher
        await _answer(sessions, sermons).execute(
            actor_account_id=uuid4(), session_id=start.session_id, attended=True
        )


# --- S-4 : la présence déclarée (le « oui » pose une présence, additive) ---


class _FakeCulte(CulteAttendancePort):
    def __init__(self):
        self.calls = []

    async def mark_declared_present(self, *, tenant_id, member_account_id, on_date, now):
        self.calls.append((tenant_id, member_account_id, on_date))


async def test_yes_marks_a_declared_presence_at_the_culte():
    tenant = uuid4()
    member, ms = _member_ms(tenant)
    sermon = _published(tenant, uuid4())
    sessions, sermons = _FakeSessions(), _FakeSermons([sermon])
    culte = _FakeCulte()
    start = await _start(sessions, sermons, ms).execute(
        actor_account_id=member, sermon_id=sermon.id
    )
    await AnswerAttendance(sessions, sermons, culte, clock=lambda: _NOW).execute(
        actor_account_id=member, session_id=start.session_id, attended=True
    )
    assert culte.calls == [(tenant, member, sermon.preached_on)]  # au culte du jour du sermon


async def test_no_marks_no_presence():
    tenant = uuid4()
    member, ms = _member_ms(tenant)
    sermon = _published(tenant, uuid4())
    sessions, sermons = _FakeSessions(), _FakeSermons([sermon])
    culte = _FakeCulte()
    start = await _start(sessions, sermons, ms).execute(
        actor_account_id=member, sermon_id=sermon.id
    )
    await AnswerAttendance(sessions, sermons, culte, clock=lambda: _NOW).execute(
        actor_account_id=member, session_id=start.session_id, attended=False
    )
    assert culte.calls == []  # le « non » n'écrit jamais de présence


async def test_companion_closes_immediately_when_a_branch_is_empty():
    tenant = uuid4()
    member, ms = _member_ms(tenant)
    sermon = _draft(tenant, uuid4())
    sermon.attach_digest(  # digest sans points essentiels → branche « non » vide
        SermonDigest(summary="s", key_points=(), capsules=(), questions=())
    )
    sermon.approve(now=_NOW)
    sermon.publish(now=_NOW)
    sessions, sermons = _FakeSessions(), _FakeSermons([sermon])
    start = await _start(sessions, sermons, ms).execute(
        actor_account_id=member, sermon_id=sermon.id
    )
    card = await _answer(sessions, sermons).execute(
        actor_account_id=member, session_id=start.session_id, attended=False
    )
    assert card.stage == "closing" and card.done


# --- S-2 : les capsules au fil (à la publication) ---


class _FakeFeed(CapsuleFeedPort):
    def __init__(self):
        self.published = []

    async def publish_capsules(self, *, tenant_id, author_account_id, capsules):
        self.published.append((tenant_id, author_account_id, list(capsules)))


async def test_publishing_pushes_the_capsules_to_the_feed():
    tenant = uuid4()
    pastor, ms = _pastor(tenant)
    s = _draft(tenant, pastor)
    s.attach_digest(_DIGEST)  # 1 capsule
    s.approve(now=_NOW)
    feed = _FakeFeed()
    await PublishSermon(_FakeSermons([s]), _access(ms), feed, clock=lambda: _NOW).execute(
        actor_account_id=pastor, sermon_id=s.id
    )
    assert feed.published
    tenant_id, author, capsules = feed.published[0]
    assert tenant_id == tenant and author == pastor and len(capsules) == 1


async def test_publishing_without_capsules_touches_no_feed():
    tenant = uuid4()
    pastor, ms = _pastor(tenant)
    s = _draft(tenant, pastor)  # aucun digest
    s.approve(now=_NOW)
    feed = _FakeFeed()
    await PublishSermon(_FakeSermons([s]), _access(ms), feed, clock=lambda: _NOW).execute(
        actor_account_id=pastor, sermon_id=s.id
    )
    assert feed.published == []


# --- S-5 : l'ingestion PDF / PPTX (un port, des adaptateurs) ---


class _FakeExtractor(SermonTextExtractor):
    def __init__(self, text="texte extrait du fichier"):
        self._text = text
        self.calls = []

    async def extract(self, data, *, kind):
        self.calls.append((data, kind))
        return self._text


async def test_uploading_a_file_extracts_text_then_deposits():
    tenant = uuid4()
    pastor, ms = _pastor(tenant)
    sermons = _FakeSermons()
    extractor = _FakeExtractor("Le fils prodigue revint vers son père.")
    cmd = DepositSermon(sermons, _access(ms), _FakeDigester(), extractor, clock=lambda: _NOW)
    dto = await cmd.execute_file(
        actor_account_id=pastor, tenant_id=tenant, title="Depuis un PDF",
        data=b"%PDF-1.4 ...", kind=SermonSourceKind.PDF, preached_on=_SUNDAY,
    )
    assert dto.source_kind == "pdf" and dto.raw_text == "Le fils prodigue revint vers son père."
    assert dto.digest is not None  # digéré comme un dépôt texte (l'IA ne voit que du texte)
    assert extractor.calls[0][1] is SermonSourceKind.PDF


async def test_uploading_without_an_extractor_is_refused():
    tenant = uuid4()
    pastor, ms = _pastor(tenant)
    cmd = DepositSermon(_FakeSermons(), _access(ms), clock=lambda: _NOW)  # pas d'extracteur
    with pytest.raises(UnsupportedSermonFormatError):
        await cmd.execute_file(
            actor_account_id=pastor, tenant_id=tenant, title="x",
            data=b"...", kind=SermonSourceKind.PDF, preached_on=_SUNDAY,
        )


async def test_a_non_member_cannot_upload_a_sermon():
    tenant, stranger = uuid4(), uuid4()
    ms = _FakeMemberships([_member(stranger, tenant)])
    cmd = DepositSermon(_FakeSermons(), _access(ms), None, _FakeExtractor(), clock=lambda: _NOW)
    with pytest.raises(UnauthorizedGroupActionError):
        await cmd.execute_file(
            actor_account_id=stranger, tenant_id=tenant, title="x",
            data=b"...", kind=SermonSourceKind.PDF, preached_on=_SUNDAY,
        )


async def test_text_extractor_decodes_bytes():
    x = CompositeTextExtractor()
    out = await x.extract("Bonjour, Église.".encode(), kind=SermonSourceKind.TEXT)
    assert out == "Bonjour, Église."


async def test_audio_is_not_supported_yet():
    x = CompositeTextExtractor()
    with pytest.raises(UnsupportedSermonFormatError):  # S-6
        await x.extract(b"...", kind=SermonSourceKind.AUDIO)


async def test_pptx_extraction_round_trip():
    from pptx import Presentation
    from pptx.util import Inches

    buf = BytesIO()
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])  # « Titre seul »
    slide.shapes.title.text = "La grâce du Père"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "Le fils prodigue revint."
    deck.save(buf)
    text = await CompositeTextExtractor().extract(buf.getvalue(), kind=SermonSourceKind.PPTX)
    assert "La grâce du Père" in text and "prodigue" in text


async def test_pdf_extraction_handles_a_real_pdf_without_crashing():
    from pypdf import PdfWriter

    buf = BytesIO()
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.write(buf)
    # Une page blanche n'a pas de texte : l'extraction rend "" sans planter (intégration pypdf OK).
    assert await CompositeTextExtractor().extract(buf.getvalue(), kind=SermonSourceKind.PDF) == ""


# --- Le sujet de reconnaissance : la seule parole du membre qui dit que ça va -------------


class _Intake:
    def __init__(self, accepted=True):
        self.submitted = []
        self._accepted = accepted

    async def submit(self, fact):
        from app.contexts.watch.application.intake import IntakeResult

        self.submitted.append(fact)
        return IntakeResult(accepted=self._accepted)


async def test_gratitude_is_deposited_for_oneself_and_never_for_another():
    """Le service ne prend **aucun identifiant de sujet** — même règle que l'anniversaire.

    Un responsable qui pourrait déposer « elle va bien » à la place de quelqu'un ferait taire un
    cas avec sa propre impression."""
    from app.contexts.sermon.application.commands.gratitude import DepositGratitude
    from app.contexts.watch.domain.facts import FactKind
    from app.contexts.watch.domain.registry import COMPANION

    intake, awa, tenant = _Intake(), uuid4(), uuid4()
    now = datetime(2026, 4, 12, tzinfo=UTC)

    accepted = await DepositGratitude(intake, clock=lambda: now).execute(
        actor_account_id=awa, tenant_id=tenant, subject="Mon fils a retrouvé du travail."
    )

    (fact,) = intake.submitted
    assert accepted is True
    assert fact.subject_id == awa  # le sujet **est** l'acteur, il n'y a pas d'autre chemin
    assert (fact.kind, fact.source) == (FactKind.GRATITUDE_DEPOSITED, COMPANION)
    assert fact.payload["subject"] == "Mon fils a retrouvé du travail."


async def test_thanking_twice_is_two_gestures():
    """Rendre grâce deux fois n'est pas un doublon : le second signe de vie vaut le premier."""
    from app.contexts.sermon.application.commands.gratitude import DepositGratitude

    intake, awa, tenant = _Intake(), uuid4(), uuid4()
    now = datetime(2026, 4, 12, tzinfo=UTC)
    command = DepositGratitude(intake, clock=lambda: now)

    await command.execute(actor_account_id=awa, tenant_id=tenant, subject="Merci.")
    await command.execute(actor_account_id=awa, tenant_id=tenant, subject="Merci encore.")

    first, second = intake.submitted
    assert first.fact_id != second.fact_id


async def test_a_disconnected_engine_never_blocks_the_gesture():
    """Ce que la personne a voulu dire ne dépend pas de l'état d'un greffon."""
    from app.contexts.sermon.application.commands.gratitude import DepositGratitude

    now = datetime(2026, 4, 12, tzinfo=UTC)
    assert await DepositGratitude(None, clock=lambda: now).execute(
        actor_account_id=uuid4(), tenant_id=uuid4(), subject="Merci."
    ) is False
