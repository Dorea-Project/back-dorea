"""Les deux écrivains — **les seuls fichiers du dépôt qui importent `pptx` et `docx`**.

Même posture que `adapters/mistral.py` : la bibliothèque tierce reste à la bordure, et l'import
est **paresseux** (patron de l'extracteur de sermon) — un dépôt sans `python-docx` installé
continue de démarrer, de résoudre des passages et de valider des citations. Seul le rendu échoue,
et il échoue en le disant.

## Ce que ces fichiers ne font pas

**Ils ne décident rien.** La frontière écran/note est déjà tenue par les types (`Deck` n'a nulle
part où mettre une mise en garde) ; ces modules mettent en page ce qu'on leur donne. Si un jour
une mise en garde apparaît sur une diapositive, ce ne sera pas ici qu'il faudra chercher —
ce sera que quelqu'un aura ajouté un champ au type.

**Ils ne produisent rien qui n'ait été validé.** C'est le service qui refuse de rendre un
livrable non `conforme` ; ici on suppose le contrôle déjà passé.

## Deux détails qui se voient depuis le fond de la salle

**Le 16:9 se pose explicitement.** `python-pptx` ouvre en 4:3 : un cadre noir de chaque côté est
la première chose que l'assemblée voit du travail de la semaine.

**La mention de destination est en pied de CHAQUE page de la note**, jamais une page de garde —
elle ne survivrait ni à une capture d'écran, ni à un partage partiel, ni à une impression recto.
"""

from __future__ import annotations

from io import BytesIO

from app.contexts.urim.deliverable.domain.documents import Deck, Note
from app.contexts.urim.deliverable.infrastructure import charte

#: 16:9 en EMU (914 400 par pouce) — 13,333 x 7,5 pouces.
_LARGEUR_16_9 = 12192000
_HAUTEUR_16_9 = 6858000

#: La géométrie de la diapositive, en EMU. Une marge d'un pouce : moins, et le texte touche le
#: bord d'un écran mal réglé — ce qui arrive dans toutes les salles.
_MARGE = 914400
_LARGEUR_UTILE = _LARGEUR_16_9 - 2 * _MARGE

#: Le filet de brique : deux points d'épaisseur, un tiers de la largeur utile. Il marque, il
#: ne souligne pas.
_FILET = 25400
_LARGEUR_FILET = _LARGEUR_UTILE // 3

#: La bande d'un pied de page ou d'un intitulé.
_PIED = 457200

#: Le texte projeté commence sous l'intitulé et s'arrête avant le pied.
_HAUT_DU_TEXTE = _MARGE + 3 * _PIED // 2
_HAUTEUR_DU_TEXTE = _HAUTEUR_16_9 - _HAUT_DU_TEXTE - _MARGE - _PIED

#: Les codes de section, rendus lisibles. Le dictionnaire ne **ferme** rien : un code inconnu
#: s'imprime tel quel, parce que la colonne est libre et qu'un pasteur peut nommer ses sections.
_SECTIONS = {
    "titre": "Titre",
    "introduction": "Introduction",
    "proposition": "Proposition",
    "phrase_interrogative": "Question",
    "phrase_de_transition": "Transition",
    "divisions": "Point",
    "subdivisions": "Sous-point",
    "illustrations": "Illustration",
    "application": "Application",
    "conclusion": "Conclusion",
    "objectif": "Objectif",
    "contexte": "Contexte",
    "definitions": "Définition",
    "nb": "NB",
    "temoignage": "Témoignage",
}

# ---------------------------------------------------------------- le vocabulaire, en clair
#
# ⚠️ **Une note de préparation n'est pas un article de revue.** « Locus », « proof-texting »,
# « pneumatologie » sont le vocabulaire des biblistes ; le pasteur qui ouvre ce document un
# vendredi soir n'a pas à le traduire pour se servir de son propre travail. Les codes du corpus
# restent ce qu'ils sont **en base** — ils sont une clé, pas une phrase — et c'est ici, au seul
# endroit qui parle à quelqu'un, qu'ils redeviennent du français.
#
# ⚠️ **Ce que cette table ne peut pas corriger** : les *motifs* de curation sont rédigés par le
# modèle et sont eux aussi techniques (« tension entre le déjà et le pas encore »). Les rendre
# lisibles est un travail d'**invite de curation**, pas de mise en page — les réécrire ici
# reviendrait à faire dire à un relecteur ce qu'il n'a pas écrit.

#: Les dix, tels qu'un prédicateur les nomme. Le libellé du corpus dit « Pneumatologie — le
#: Saint-Esprit » ; on garde la moitié qui parle.
_LOCI = {
    "theologie_propre": "Dieu lui-même",
    "christologie": "Jésus-Christ",
    "pneumatologie": "le Saint-Esprit",
    "anthropologie": "l'homme",
    "hamartiologie": "le péché",
    "soteriologie": "le salut",
    "ecclesiologie": "l'Église",
    "angelologie": "les anges",
    "demonologie": "Satan et les démons",
    "eschatologie": "les derniers temps",
}

#: Les quatre forces. `resiste` garde son avertissement : c'est celle qui protège.
_FORCES = {
    "dominant": "au cœur du texte",
    "porte": "présent, en appui",
    "resiste": "⚠ complique ce point",
    "absent": "le texte n'en dit rien",
}

#: « proof-texting » ne se traduit pas, il s'explique : c'est faire dire au texte ce qu'on
#: voulait déjà entendre.
_RISQUES = {
    "faible": "peu de risque de faire dire au texte plus qu'il ne dit",
    "moyen": "attention à ne pas faire dire au texte plus qu'il ne dit",
    "eleve": "⚠ risque réel de faire dire au texte ce qu'on voulait déjà entendre",
}

_PLANS = {
    "thematique": "un plan par thème",
    "expositif": "un plan verset par verset",
    "textuel": "un plan collé au texte",
}

_MATIERES = {
    "doctrinal": "une doctrine",
    "ethique": "une conduite",
    "biographique": "un personnage",
    "historique": "un récit",
    "typologique": "une figure",
    "prophetique": "une annonce",
}


def _clair(valeur: str, table: dict[str, str]) -> str:
    """Le mot du corpus, rendu en français — **et tel quel si on ne le connaît pas**.

    Un code inconnu s'affiche plutôt que de disparaître : mieux vaut un mot technique qu'un
    trou dans la note de quelqu'un."""
    if valeur in table:
        return table[valeur]
    # Le corpus écrit parfois « Pneumatologie — le Saint-Esprit » : on garde la moitié droite.
    return valeur.split("—")[-1].strip() if "—" in valeur else valeur


#: Ce qu'un pied de page de note doit dire, et à qui elle est destinée. Voir l'en-tête.
MENTION = (
    "Note de préparation — les mises en garde s'adressent au prédicateur, "
    "non à l'assemblée."
)


class RenduIndisponibleError(RuntimeError):
    """La bibliothèque de rendu n'est pas installée.

    Une erreur nommée plutôt qu'un `ImportError` nu : sur un déploiement incomplet, la phrase
    doit dire **quoi faire**, comme `CorpusNonSemeError` le fait pour le corpus."""


def rendre_deck(deck: Deck) -> bytes:
    """Le `.pptx` — une couverture, puis une diapositive par texte projeté.

    ⚠️ **Rien ne vient du gabarit par défaut.** `python-pptx` ouvre sur un modèle Office —
    fond blanc, Calibri, filets bleus — et une présentation qui garde ces traits n'est celle
    de personne. On dessine donc sur la disposition vide et on pose tout : le fond, le filet,
    la référence, le pied de page.
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Emu, Pt
    except ImportError as exc:  # pragma: no cover - dépend de l'installation
        raise RenduIndisponibleError(
            "`python-pptx` n'est pas installé : le rendu des diapositives est indisponible."
        ) from exc

    presentation = Presentation()
    presentation.slide_width = _LARGEUR_16_9
    presentation.slide_height = _HAUTEUR_16_9

    #: La disposition **vide** — la septième du gabarit. Les autres portent des cadres de
    #: texte qu'on ne veut pas : un « Cliquez pour ajouter un titre » resté en place au mur
    #: d'une assemblée, ça arrive, et ça ne s'oublie pas.
    vierge = presentation.slide_layouts[6]

    def zone(page, gauche, haut, largeur, hauteur):
        cadre = page.shapes.add_textbox(
            Emu(gauche), Emu(haut), Emu(largeur), Emu(hauteur)
        ).text_frame
        cadre.word_wrap = True
        return cadre

    def ecrire(cadre, texte, *, taille, couleur, police, gras=False, premier=False):
        paragraphe = cadre.paragraphs[0] if premier else cadre.add_paragraph()
        paragraphe.text = texte
        for morceau in paragraphe.runs:
            morceau.font.size = Pt(taille)
            morceau.font.color.rgb = RGBColor(*couleur)
            morceau.font.name = police
            morceau.font.bold = gras
        return paragraphe

    def fond(page, couleur):
        page.background.fill.solid()
        page.background.fill.fore_color.rgb = RGBColor(*couleur)

    def filet(page, haut, largeur):
        """Le trait de brique. Un rectangle sans contour — la forme la plus sûre qui soit."""
        from pptx.enum.shapes import MSO_SHAPE

        trait = page.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(_MARGE), Emu(haut), Emu(largeur), Emu(_FILET)
        )
        trait.fill.solid()
        trait.fill.fore_color.rgb = RGBColor(*charte.BRIQUE)
        trait.line.fill.background()
        trait.shadow.inherit = False

    # ------------------------------------------------------------------ la couverture
    #
    # Fond marine et titre clair : c'est la seule diapositive que l'assemblée regarde
    # pendant qu'on s'installe, et c'est là que le document dit d'où il vient.
    couverture = presentation.slides.add_slide(vierge)
    fond(couverture, charte.MARINE)
    filet(couverture, _HAUTEUR_16_9 // 3 - 2 * _PIED, _LARGEUR_FILET)
    titre = zone(couverture, _MARGE, _HAUTEUR_16_9 // 3, _LARGEUR_UTILE, _HAUTEUR_16_9 // 3)
    ecrire(
        titre,
        deck.titre,
        taille=40,
        couleur=(0xF7, 0xF4, 0xE4),
        police=charte.SERIF,
        premier=True,
    )
    # Le sous-titre reste **vide** : y mettre le thème proposé par le moteur ferait monter à
    # l'écran une phrase que personne n'a écrite.
    signature_couverture = zone(
        couverture, _MARGE, _HAUTEUR_16_9 - _MARGE - _PIED, _LARGEUR_UTILE, _PIED
    )
    ecrire(
        signature_couverture,
        charte.SIGNATURE,
        taille=11,
        couleur=charte.GRIS,
        police=charte.GROTESQUE,
        premier=True,
    )

    # ------------------------------------------------------------------ les textes projetés
    for rang, diapositive in enumerate(deck.diapositives, start=1):
        page = presentation.slides.add_slide(vierge)
        fond(page, charte.SABLE)
        filet(page, _MARGE, _LARGEUR_FILET)

        entete = zone(page, _MARGE, _MARGE + _FILET * 3, _LARGEUR_UTILE, _PIED)
        ecrire(
            entete,
            (diapositive.titre or diapositive.reference).upper(),
            taille=14,
            couleur=charte.BRIQUE,
            police=charte.GROTESQUE,
            gras=True,
            premier=True,
        )

        corps = zone(page, _MARGE, _HAUT_DU_TEXTE, _LARGEUR_UTILE, _HAUTEUR_DU_TEXTE)
        ecrire(
            corps,
            diapositive.texte_projete,
            taille=charte.corps_du_verset(diapositive.texte_projete),
            couleur=charte.ENCRE,
            police=charte.SERIF,
            premier=True,
        )

        pied = zone(page, _MARGE, _HAUTEUR_16_9 - _MARGE - _PIED, _LARGEUR_UTILE, _PIED)
        # La référence en pied — une projection sans référence est une citation invérifiable
        # pour qui la lit depuis le banc. ⚠️ **Sauf quand elle est déjà l'intitulé** : sans ce
        # garde, une diapositive sans titre affiche deux fois la même référence, ce qui se voit
        # du fond de la salle et fait amateur.
        rappel = diapositive.reference if diapositive.titre else charte.SIGNATURE
        ecrire(
            pied,
            f"{rappel}   ·   {rang}",
            taille=11,
            couleur=charte.GRIS,
            police=charte.GROTESQUE,
            premier=True,
        )

    flux = BytesIO()
    presentation.save(flux)
    return flux.getvalue()


def rendre_note(note: Note) -> bytes:
    """Le `.docx` — tout ce que le moteur a rassemblé, et rien qui monte à l'écran."""
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError as exc:  # pragma: no cover - dépend de l'installation
        raise RenduIndisponibleError(
            "`python-docx` n'est pas installé : le rendu de la note est indisponible."
        ) from exc

    document = Document()
    _habiller(document)
    _pied_de_page(document)

    _titre_du_document(document, note.titre or note.reference, note.reference)

    if note.unite:
        document.add_heading("L'unité littéraire", level=1)
        document.add_paragraph(note.unite)
        # ⚠️ **Le motif du découpage voyage avec les bornes.** C'est la phrase que le pasteur
        # lit pour vous contredire ; l'imprimer sans elle donnerait des bornes venues de nulle
        # part.
        _sous_texte(document, note.motif_unite)

    # ⚠️ **La section existe meme vide, et c'est le contraire d'un oubli.**
    #
    # Une section absente se lit comme une negligence du document ; une section qui
    # dit « a ecrire » se lit comme une place laissee. La note sans plan est un
    # document de travail — tout ce que le moteur a etabli, et un cadre pour ce que
    # le pasteur va ecrire dessus.
    document.add_heading("Votre plan", level=1)
    if not note.plan:
        _sous_texte(
            document,
            "À écrire. Le document met en page ce que vous aurez écrit ; il ne "
            "l'écrit pas à votre place.",
        )

    if note.plan:
        for code, corps, appuis in note.plan:
            # ⚠️ **Chaque point est un TITRE, pas une puce.** En liste, les trois points d'un
            # sermon tiennent en cinq lignes et le document n'offre nulle part où les
            # développer — or c'est là que le travail se fait. Un titre ouvre la place ;
            # une puce la ferme.
            document.add_heading(corps, level=2)
            _sous_texte(document, _SECTIONS.get(code, code))
            # **Ce qui développe le point sans l'écrire** : les textes que le pasteur a
            # lui-même convoqués dans sa ligne, servis dessous. Urim n'ajoute pas une phrase
            # de sermon — il pose sous la main ce qu'il faudrait aller chercher.
            for reference, servi in appuis:
                appui = document.add_paragraph(style="List Bullet")
                appui.add_run(f"{reference} — ").bold = True
                appui.add_run(servi)

    if note.versets:
        document.add_heading("Le texte", level=1)
        # ⚠️ **Dire dans quelle version il a préparé.** Sur Romains 8:1, l'Ostervald porte
        # une clause que la LSG omet, et les deux versets font deux sermons opposés. Une
        # note muette sur sa source laisse croire qu'il n'y en avait qu'une.
        if note.version:
            _sous_texte(document, f"Version préparée : {note.version}")
        for reference, corps in note.versets:
            document.add_paragraph(f"{reference} — {corps}")

    if note.pesees:
        document.add_heading("Ce dont ce texte parle", level=1)
        _pesees(document, note)

    if note.mises_en_garde:
        document.add_heading("Ce que ce texte ne dit pas", level=1)
        for garde in note.mises_en_garde:
            document.add_paragraph(garde, style="List Bullet")

    if note.resistances:
        document.add_heading("Ailleurs, des textes en tension avec le vôtre", level=1)
        # **La rubrique qui protège du proof-texting**, et la seule que le pasteur ne
        # trouvera pas seul : ce sont précisément les textes qu'on ne cherche pas quand
        # on a déjà son idée. Ils viennent d'AUTRES livres, sur l'axe qu'il a retenu.
        _sous_texte(
            document,
            "Ces passages compliquent l'axe que vous avez retenu. Les rencontrer avant "
            "dimanche vaut mieux que de les entendre après.",
        )
        for reference, motif in note.resistances:
            ligne = document.add_paragraph(style="List Bullet")
            ligne.add_run(reference).bold = True
            _sous_texte(document, motif)

    if note.faisabilites:
        document.add_heading("Quels plans tiennent sur ce texte", level=1)
        # Les faisables d'abord : six refus d'affilée se lisent comme une liste de portes
        # fermées, alors que l'information utile est **par où passer**.
        for couple, faisable, refus, risque in sorted(
            note.faisabilites, key=lambda f: not f[1]
        ):
            libelle = " sur ".join(
                _clair(part.strip(), table)
                for part, table in zip(couple.split(" x "), (_PLANS, _MATIERES), strict=False)
            )
            verdict = "tient" if faisable else f"ne tient pas — {refus}"
            document.add_paragraph(f"{libelle} : {verdict}", style="List Bullet")
            if risque and faisable:
                _sous_texte(document, _clair(risque, _RISQUES))

    if note.appuis:
        document.add_heading("Les textes que vous convoquez", level=1)
        for reference, corps, verdict in note.appuis:
            # ⚠️ **Une saisie illisible s'imprime avec son motif**, jamais en silence : la
            # perdre obligerait le pasteur à se souvenir de ce qu'il voulait citer.
            document.add_paragraph(f"{reference or verdict} — {corps}", style="List Bullet")

    if note.original:
        document.add_heading("Les mots d'origine, et où ils reviennent", level=1)
        # ⚠️ **Urim ne traduit pas, et le dit.** MorphGNT ne porte aucune traduction et les
        # lexiques libres sont en anglais : une glose produite par un modèle aurait l'air d'une
        # source, et personne ne relit une définition grecque avant de la redire en chaire.
        # Ce qui la remplace est plus sûr et souvent plus parlant — les autres endroits où le
        # mot paraît. La culture d'un mot s'enseigne par sa récurrence, pas par un synonyme.
        # ⚠️ **L'attribution est une obligation de licence, pas une politesse** : le sens
        # vient de TBESG (STEPBible, CC BY 4.0), traduit. Elle voyage donc avec la section.
        _sous_texte(
            document,
            "Le sens est traduit d'un lexique publié (STEP Bible, CC BY) et l'entrée "
            "d'origine est donnée à côté. Les autres passages où le mot paraît disent le "
            "reste — c'est l'usage qui l'enseigne.",
        )
        for mot in note.original:
            ligne = document.add_paragraph(style="List Bullet")
            ligne.add_run(mot.forme).bold = True
            # **La phonétique d'abord** : sans elle, le mot reste étranger même sous les yeux
            # de celui qui le lit — il ne peut ni le dire, ni le retenir, ni le rechercher.
            if mot.phonetique:
                prononce = ligne.add_run(f"  [{mot.phonetique}]")
                prononce.italic = True
            detail = f" ({mot.lemme})" if mot.lemme and mot.lemme != mot.forme else ""
            grammaire = " · ".join(p for p in (mot.nature, mot.morphologie) if p)
            ligne.add_run(
                f"{detail} — {mot.reference}{' · ' + grammaire if grammaire else ''}"
            )
            # **Le sens, et sa source dans la même ligne.** Les séparer les rendrait
            # indépendants : la traduction pourrait circuler sans ce qui la vérifie.
            if mot.sens:
                sens = document.add_paragraph()
                fort = sens.add_run(mot.sens)
                fort.bold = True
                appui = sens.add_run(
                    f"  ({mot.sens_source}"
                    + (f", {mot.strong}" if mot.strong else "")
                    + ")"
                )
                appui.italic = True
                appui.font.size = Pt(9)
            # ⚠️ **Ce que ces versets ont en commun**, présenté comme un fait et jamais comme
            # une définition : le mot français qui revient là où le mot grec paraît est presque
            # toujours sa traduction — et quand il ne l'est pas, le pasteur le voit, parce que
            # les versets sont juste dessous.
            if mot.communs:
                _sous_texte(
                    document,
                    "ces versets ont en commun : " + ", ".join(mot.communs),
                )
            # ⚠️ **Le verset entier, pas seulement sa référence.** « revient en Luc 15:22 » ne
            # dit rien à personne ; « on lui mit des souliers aux pieds » dit ce qu'est
            # l'objet. C'est la demande d'un sens littéral, servie par le texte lui-même.
            for ou, servi in mot.ailleurs:
                _sous_texte(document, f"{ou} — {servi}")

    if note.ecartees:
        document.add_heading("Ce que vous avez écarté en chemin", level=1)
        for option, motif in note.ecartees:
            document.add_paragraph(option, style="List Bullet")
            _sous_texte(document, motif)

    provenance = document.add_paragraph()
    trace = provenance.add_run(
        f"Relecture de ce passage : {note.signature or 'aucune'} · "
        f"état du corpus {note.corpus_snapshot or 'inconnu'}"
    )
    trace.font.size = Pt(8)

    flux = BytesIO()
    document.save(flux)
    return flux.getvalue()


def _nom_de_locus(axe: str) -> str:
    """**Le nom savant ET ce qu'il désigne** — « Christologie — Jésus-Christ ».

    ⚠️ Le remplacer par sa seule glose lui fait perdre ce qu'il est : les dix loci sont un
    vocabulaire fixe, celui dans lequel un pasteur a été formé et dans lequel il retrouvera son
    travail d'une semaine à l'autre. Traduire seul efface la clé ; laisser seul le terme
    technique exclut celui qui ne l'a pas appris. Les deux, donc — le savant d'abord, parce
    que c'est lui qui nomme."""
    clair = _clair(axe, _LOCI)
    savant = axe.split("—")[0].strip().replace("_", " ").capitalize()
    return f"{savant} — {clair}" if clair.lower() != savant.lower() else savant


def _pesees(document, note: Note) -> None:
    """Ce que le texte porte — **et qui le dit**.

    Trois régimes, parce que ces lignes n'ont pas le même statut :

    - **son choix** et **le dominant** sont développés : ce sont eux qui décident du sermon,
      et ils **ne coïncident pas toujours**. Nommer les deux côte à côte est l'information
      la plus utile de la page ;
    - ce qui **porte** ou **résiste** est nommé avec son motif, sans être mis en avant ;
    - ce dont le texte **ne dit rien** tient en une ligne, à la fin. La distinction reste
      (*quelqu'un a regardé* ≠ *personne n'a regardé*), mais quatre paragraphes pour dire
      « rien ici » repoussaient en page 3 ce que le pasteur vient chercher.
    """
    from docx.shared import Pt

    retenu = [p for p in note.pesees if note.axe_retenu and p[0] == note.axe_retenu]
    dominants = [p for p in note.pesees if p[1] == "dominant" and p not in retenu]
    autres = [
        p for p in note.pesees
        if p not in retenu and p not in dominants and p[1] != "absent"
    ]
    muets = [p for p in note.pesees if p[1] == "absent" and p not in retenu]

    for pesees, etiquette in ((retenu, "votre choix"), (dominants, "le dominant du texte")):
        for axe, force, motif in pesees:
            titre = document.add_paragraph()
            titre.add_run(_nom_de_locus(axe)).bold = True
            marque = titre.add_run(f"  ({etiquette})")
            marque.italic = True
            marque.font.size = Pt(9)
            if etiquette == "votre choix" and force != "dominant":
                # Le désaccord se dit, il ne se corrige pas : il a peut-être raison, et c'est
                # lui qui prêche.
                _sous_texte(
                    document,
                    f"Le corpus le classe « {_clair(force, _FORCES)} » sur cette unité.",
                )
            document.add_paragraph(motif)

    for axe, force, motif in autres:
        ligne = document.add_paragraph(style="List Bullet")
        ligne.add_run(f"{_nom_de_locus(axe)} — {_clair(force, _FORCES)}")
        _sous_texte(document, motif)

    if muets:
        _sous_texte(
            document,
            "Relus et sans objet ici : "
            + ", ".join(_clair(axe, _LOCI) for axe, _f, _m in muets)
            + ". (Relus, pas oubliés : le texte n'en dit rien.)",
        )


def _habiller(document) -> None:
    """La charte, posée sur les styles du document — **une fois, pas paragraphe par paragraphe**.

    Word applique ses propres styles : Calibri 11, titres bleu Office. Une note qui les garde
    est une note de Word, pas d'Urim. On réécrit donc le style « Normal » et les trois niveaux
    de titre, et tout ce qui s'ajoute ensuite en hérite.
    """
    from docx.shared import Pt, RGBColor

    normal = document.styles["Normal"]
    normal.font.name = charte.SERIF
    normal.font.size = Pt(11)
    normal.font.color.rgb = RGBColor(*charte.ENCRE)
    normal.paragraph_format.space_after = Pt(8)
    normal.paragraph_format.line_spacing = 1.25

    for niveau, taille, couleur in (
        (1, 16, charte.MARINE),
        (2, 13, charte.BRIQUE),
        (3, 11, charte.MARINE),
    ):
        style = document.styles[f"Heading {niveau}"]
        style.font.name = charte.GROTESQUE
        style.font.size = Pt(taille)
        style.font.bold = True
        style.font.color.rgb = RGBColor(*couleur)
        style.paragraph_format.space_before = Pt(16)
        style.paragraph_format.space_after = Pt(4)


def _titre_du_document(document, titre: str, reference: str) -> None:
    """Le bloc de tête : le titre, la référence sous lui, un filet.

    ⚠️ **Pas de style « Title » de Word.** Il porte une bordure bleue et une police que
    personne n'a choisie ; on écrit le bloc à la main pour qu'il soit celui de la charte.
    """
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor

    paragraphe = document.add_paragraph()
    paragraphe.paragraph_format.space_after = Pt(2)
    morceau = paragraphe.add_run(titre)
    morceau.font.size = Pt(24)
    morceau.font.name = charte.SERIF
    morceau.font.bold = True
    morceau.font.color.rgb = RGBColor(*charte.MARINE)

    sous_titre = document.add_paragraph()
    sous_titre.paragraph_format.space_after = Pt(10)
    rappel = sous_titre.add_run(reference.upper())
    rappel.font.size = Pt(10)
    rappel.font.name = charte.GROTESQUE
    rappel.font.bold = True
    rappel.font.color.rgb = RGBColor(*charte.BRIQUE)

    # Le filet : une bordure basse sur un paragraphe vide. C'est la seule façon de tracer un
    # trait en `python-docx` sans dessiner une image.
    _filet(document.add_paragraph())
    document.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.LEFT


def _filet(paragraphe) -> None:
    """Un trait de brique sous un paragraphe vide."""
    from docx.oxml.ns import qn
    from docx.oxml.parser import OxmlElement

    bordures = OxmlElement("w:pBdr")
    bas = OxmlElement("w:bottom")
    bas.set(qn("w:val"), "single")
    bas.set(qn("w:sz"), "12")
    bas.set(qn("w:space"), "1")
    bas.set(qn("w:color"), "CC3C1F")
    bordures.append(bas)
    paragraphe._p.get_or_add_pPr().append(bordures)


def _sous_texte(document, texte: str) -> None:
    if not texte:
        return
    from docx.shared import Pt, RGBColor

    paragraphe = document.add_paragraph()
    morceau = paragraphe.add_run(texte)
    morceau.italic = True
    morceau.font.size = Pt(9)
    morceau.font.name = charte.GROTESQUE
    morceau.font.color.rgb = RGBColor(*charte.GRIS)


def _pied_de_page(document) -> None:
    """La mention de destination, **dans le pied de page de chaque page**.

    Une page de garde ne survit ni à une capture d'écran, ni à un partage partiel, ni à une
    impression recto — et le PDF, lui, circule pour de bon."""
    from docx.shared import Pt, RGBColor

    for section in document.sections:
        paragraphe = section.footer.paragraphs[0]
        morceau = paragraphe.add_run(MENTION)
        morceau.italic = True
        morceau.font.size = Pt(8)
        morceau.font.name = charte.GROTESQUE
        morceau.font.color.rgb = RGBColor(*charte.GRIS)

        # La signature suit la mention, sur la même ligne : deux lignes de pied mangeraient une
        # ligne de travail sur chaque page.
        signature = paragraphe.add_run(f"   ·   {charte.SIGNATURE}")
        signature.font.size = Pt(8)
        signature.font.name = charte.GROTESQUE
        signature.font.color.rgb = RGBColor(*charte.GRIS)
