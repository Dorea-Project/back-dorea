"""Extracteurs de texte de sermon (S-5) — un fichier déposé → du **texte** pour la digestion.

Un seul port (`SermonTextExtractor`), plusieurs adaptateurs derrière un **dispatcher** par format :
- texte : décodage UTF-8 ;
- **PDF** : `pypdf` (import paresseux) ;
- **PPTX** : `python-pptx` (import paresseux) ;
- **audio** : non pris en charge ici (S-6) → erreur claire.

L'IA (S-1) ne voit jamais que le texte qui sort d'ici. Ajouter un format = ajouter un adaptateur,
sans toucher au domaine ni au pipeline.
"""

from __future__ import annotations

import zipfile
from io import BytesIO

from app.contexts.sermon.application.ports import SermonTextExtractor
from app.contexts.sermon.domain.enums import SermonSourceKind
from app.contexts.sermon.domain.errors import (
    SermonFileTooComplexError,
    UnsupportedSermonFormatError,
)

# --- Garde-fous DOREA-011 : le plafond d'upload borne ce qui ENTRE (15 Mo),
# ceux-ci bornent ce que le fichier COÛTE à ouvrir.
_MAX_UNCOMPRESSED = 200 * 1024 * 1024  # 200 Mio décompressés, tous membres confondus
_MAX_RATIO = 200  # taux d'expansion : au-delà, c'est une bombe, pas un sermon
_MAX_PDF_PAGES = 2_000  # un sermon n'a pas deux mille pages
_MAX_TEXT_CHARS = 400_000  # borne du texte rendu (aussi DOREA-026 : entrée du modèle)
_CHUNK = 64 * 1024


def _capped(text: str) -> str:
    """Un sermon plus long que ça est une anomalie — on tronque plutôt que d'engorger."""
    return text[:_MAX_TEXT_CHARS]


def _extract_text(data: bytes) -> str:
    return _capped(data.decode("utf-8", errors="replace").strip())


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader  # import paresseux (dépendance de S-5)

    reader = PdfReader(BytesIO(data))
    if len(reader.pages) > _MAX_PDF_PAGES:
        raise SermonFileTooComplexError(
            "Ce PDF compte trop de pages pour être dépouillé.",
            details={"pages": len(reader.pages), "max": _MAX_PDF_PAGES},
        )
    pages: list[str] = []
    total = 0
    for page in reader.pages:
        text = (page.extract_text() or "").strip()
        if not text:
            continue
        pages.append(text)
        total += len(text)
        if total >= _MAX_TEXT_CHARS:  # on s'arrête net, sans lire le reste
            break
    return _capped("\n".join(pages).strip())


def _guard_zip_bomb(data: bytes) -> None:
    """Un `.pptx` est un ZIP : on mesure ce qu'il pèse **décompressé** avant de l'ouvrir.

    On ne se fie pas à la taille déclarée dans l'en-tête (elle ment quand on l'attaque) :
    on décompresse réellement, par tranches, en s'arrêtant dès le plafond franchi."""
    try:
        archive = zipfile.ZipFile(BytesIO(data))
    except zipfile.BadZipFile as exc:
        raise UnsupportedSermonFormatError("Fichier PPTX illisible.") from exc

    total = 0
    with archive:
        for info in archive.infolist():
            with archive.open(info) as member:
                while chunk := member.read(_CHUNK):
                    total += len(chunk)
                    if total > _MAX_UNCOMPRESSED or total > len(data) * _MAX_RATIO:
                        raise SermonFileTooComplexError(
                            "Ce fichier se déploie démesurément une fois décompressé.",
                            details={"decompresse_octets": total, "compresse": len(data)},
                        )


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation  # import paresseux (dépendance de S-5)

    _guard_zip_bomb(data)  # avant d'ouvrir : on refuse ce qui explose
    deck = Presentation(BytesIO(data))
    lines: list[str] = []
    total = 0
    for slide in deck.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    lines.append(text)
                    total += len(text)
        if total >= _MAX_TEXT_CHARS:
            break
    return _capped("\n".join(lines).strip())


class CompositeTextExtractor(SermonTextExtractor):
    """Aiguille vers l'extracteur du format ; l'audio (S-6) n'est pas encore branché."""

    async def extract(self, data: bytes, *, kind: SermonSourceKind) -> str:
        if kind is SermonSourceKind.TEXT:
            return _extract_text(data)
        if kind is SermonSourceKind.PDF:
            return _extract_pdf(data)
        if kind is SermonSourceKind.PPTX:
            return _extract_pptx(data)
        raise UnsupportedSermonFormatError(
            "Ce format de sermon n'est pas encore pris en charge.",
            details={"kind": kind.value},
        )


def build_text_extractor() -> SermonTextExtractor:
    return CompositeTextExtractor()
